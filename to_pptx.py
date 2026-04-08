"""templates.json + design.json → 一つの PPTX (python-pptx)

使い方:
  python to_pptx.py              # すべて → all_templates.pptx
  python to_pptx.py cover        # 特定テンプレートのみ → cover.pptx
"""
import json
import os
import re
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree as _etree

_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
def _el(tag, attrib=None):
    return _etree.Element(f"{{{_NS}}}{tag}", attrib=attrib or {})

HERE = os.path.dirname(os.path.abspath(__file__))
_arg = sys.argv[1] if len(sys.argv) > 1 else None

# ── JSON 読み込み ──────────────────────────────────────
with open(os.path.join(HERE, "templates.json"), encoding="utf-8") as f:
    TEMPLATES = json.load(f)
with open(os.path.join(HERE, "design.json"), encoding="utf-8") as f:
    DESIGN = json.load(f)

C = DESIGN["COLORS"]
F = DESIGN["FONTS"]
L = DESIGN["LAYOUT"]


# ── 描画ヘルパー ───────────────────────────────────────

def rgb(hex6):
    return RGBColor.from_string(hex6)


def _no_shadow(shape):
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))


def rect(slide, x, y, w, h, fill, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(line_w or 0.75)
    else:
        shape.line.fill.background()
    _no_shadow(shape)
    return shape


def _is_cjk(ch):
    cp = ord(ch)
    return (0x3000 <= cp <= 0x9FFF or 0xF900 <= cp <= 0xFAFF or
            0xFF00 <= cp <= 0xFFEF or 0x20000 <= cp <= 0x2A6DF)


def _split_script(txt):
    if not txt:
        return []
    segs, cur_jp, cur = [], _is_cjk(txt[0]), txt[0]
    for ch in txt[1:]:
        jp = _is_cjk(ch)
        if jp == cur_jp:
            cur += ch
        else:
            segs.append((cur_jp, cur))
            cur_jp, cur = jp, ch
    segs.append((cur_jp, cur))
    return segs


def _add_runs(p, txt, size, bold, color_hex):
    sz_val = str(int(size * 100))
    b_val  = "1" if bold else "0"
    clr    = color_hex.upper()
    for is_jp, seg in _split_script(txt):
        r   = _el("r")
        rPr = _el("rPr", {"lang": "ja-JP" if is_jp else "en-US",
                           "sz": sz_val, "b": b_val, "dirty": "0"})
        fill = _el("solidFill")
        fill.append(_el("srgbClr", {"val": clr}))
        rPr.append(fill)
        rPr.append(_el("ea" if is_jp else "latin",
                       {"typeface": "Meiryo" if is_jp else "Segoe UI"}))
        t = _el("t")
        t.text = seg
        r.append(rPr); r.append(t)
        p._p.append(r)


def _fix_bodyPr(tf):
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    for old in bodyPr.findall(qn("a:spAutoFit")):
        bodyPr.remove(old)
    if not bodyPr.findall(qn("a:noAutofit")):
        bodyPr.append(bodyPr.makeelement(qn("a:noAutofit"), {}))


def _set_line_spacing(p, size_pt, mult=1.0):
    pPr = p._p.get_or_add_pPr()
    for old in pPr.findall(qn("a:lnSpc")):
        pPr.remove(old)
    lnSpc = _el("lnSpc")
    lnSpc.append(_el("spcPts", {"val": str(int(size_pt * mult * 100))}))
    pPr.append(lnSpc)


def text(slide, txt, x, y, w, h, size, bold=False, color="333333",
         anchor=MSO_ANCHOR.MIDDLE, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    _fix_bodyPr(tf)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    _add_inline_runs(p, txt, size, color, base_bold=bold)
    _no_shadow(box)
    return box


def md_content(slide, items, x, y, w, h, size, color="333333", head_size=None):
    if head_size is None:
        head_size = size
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    for old in bodyPr.findall(qn("a:spAutoFit")):
        bodyPr.remove(old)
    bodyPr.append(bodyPr.makeelement(qn("a:noAutofit"), {}))
    for idx, (item_type, item_text) in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if item_type == "heading":
            p.space_before = Pt(0) if idx == 0 else Pt(8)
            p.space_after  = Pt(3)
            item_size = head_size
            is_bold, is_italic = True, False
            eff_color = color
        elif item_type == "blockquote":
            p.space_before = Pt(2)
            p.space_after  = Pt(2)
            item_size = size
            is_bold, is_italic = False, True
            eff_color = C.get("textMuted", "888888")
        else:
            p.space_before = Pt(0)
            p.space_after  = Pt(4)
            item_size = size
            is_bold, is_italic = False, False
            eff_color = color
        pPr = p._p.get_or_add_pPr()
        for tag in (qn("a:buChar"), qn("a:buNone")):
            for old in pPr.findall(tag):
                pPr.remove(old)
        for attr in ("marL", "indent"):
            pPr.attrib.pop(attr, None)
        if item_type == "bullet":
            pPr.set("marL", "228600"); pPr.set("indent", "-228600")
            pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "\u2022"}))
        elif item_type == "bullet2":
            pPr.set("marL", "457200"); pPr.set("indent", "-228600")
            pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "\u25e6"}))
        elif item_type in ("ordered", "ordered2"):
            pPr.set("marL", "457200" if item_type == "ordered2" else "228600")
            pPr.set("indent", "-228600")
            pPr.append(pPr.makeelement(qn("a:buNone"), {}))
        elif item_type == "blockquote":
            pPr.set("marL", "228600")
            pPr.append(pPr.makeelement(qn("a:buNone"), {}))
        else:
            pPr.append(pPr.makeelement(qn("a:buNone"), {}))
        _add_inline_runs(p, item_text, item_size, eff_color, is_bold, is_italic)
        _set_line_spacing(p, item_size, mult=1.4 if item_type == "heading" else 1.8)
    _no_shadow(box)
    return box


def hline(slide, x, y, w, color, width_pt=1):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(width_pt))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    _no_shadow(shape)
    return shape


def vline(slide, x, y, h, color, width_pt=1):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Pt(width_pt), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    _no_shadow(shape)
    return shape


def down_triangle_shape(slide, x, y, w, h, fill_hex):
    """下向き三角形シェイプ（カスタムジオメトリ）"""
    from pptx.oxml import parse_xml
    x_emu = int(Inches(x)); y_emu = int(Inches(y))
    w_emu = int(Inches(w)); h_emu = int(Inches(h))
    fill  = fill_hex.upper()
    ids   = [int(e.get("id")) for e in slide.shapes._spTree.iter()
             if e.get("id") and e.get("id").isdigit()]
    shape_id = (max(ids) + 1) if ids else 100
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sp_xml = (
        f'<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
        f'<p:nvSpPr><p:cNvPr id="{shape_id}" name="Arrow{shape_id}"/>'
        '<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr>'
        f'<a:xfrm><a:off x="{x_emu}" y="{y_emu}"/>'
        f'<a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>'
        '<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        '<a:rect l="0" t="0" r="r" b="b"/>'
        '<a:pathLst><a:path w="100000" h="100000">'
        '<a:moveTo><a:pt x="0" y="0"/></a:moveTo>'
        '<a:lnTo><a:pt x="100000" y="0"/></a:lnTo>'
        '<a:lnTo><a:pt x="50000" y="100000"/></a:lnTo>'
        '<a:close/></a:path></a:pathLst></a:custGeom>'
        f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>'
        '<a:ln><a:noFill/></a:ln>'
        '<a:effectLst/></p:spPr>'
        '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    )
    slide.shapes._spTree.append(parse_xml(sp_xml))


def arrow_shape(slide, x, y, w, h, fill_hex, border_hex="CCCCCC", border_w=0.75):
    from pptx.oxml import parse_xml
    x_emu = int(Inches(x)); y_emu = int(Inches(y))
    w_emu = int(Inches(w)); h_emu = int(Inches(h))
    fill   = fill_hex.upper(); border = border_hex.upper()
    border_emu = int(border_w * 12700)
    ids = [int(e.get("id")) for e in slide.shapes._spTree.iter()
           if e.get("id") and e.get("id").isdigit()]
    shape_id = (max(ids) + 1) if ids else 100
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sp_xml = (
        f'<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
        f'<p:nvSpPr><p:cNvPr id="{shape_id}" name="Step{shape_id}"/>'
        '<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr>'
        f'<a:xfrm><a:off x="{x_emu}" y="{y_emu}"/>'
        f'<a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>'
        '<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        '<a:rect l="0" t="0" r="r" b="b"/>'
        '<a:pathLst><a:path w="100000" h="100000">'
        '<a:moveTo><a:pt x="0" y="0"/></a:moveTo>'
        '<a:lnTo><a:pt x="88000" y="0"/></a:lnTo>'
        '<a:lnTo><a:pt x="100000" y="50000"/></a:lnTo>'
        '<a:lnTo><a:pt x="88000" y="100000"/></a:lnTo>'
        '<a:lnTo><a:pt x="0" y="100000"/></a:lnTo>'
        '<a:close/></a:path></a:pathLst></a:custGeom>'
        f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>'
        f'<a:ln w="{border_emu}"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill></a:ln>'
        '<a:effectLst/></p:spPr>'
        '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    )
    slide.shapes._spTree.append(parse_xml(sp_xml))


_INLINE_RE = re.compile(
    r'\*\*([^*\n]+?)\*\*'
    r'|\*([^*\n]+?)\*'
    r'|~~([^~\n]+?)~~'
    r'|`([^`\n]+?)`'
)


def _parse_inline(txt):
    """Return list of (seg, bold, italic, strike, code)"""
    segs, last = [], 0
    for m in _INLINE_RE.finditer(txt):
        if m.start() > last:
            segs.append((txt[last:m.start()], False, False, False, False))
        if m.group(1) is not None:
            segs.append((m.group(1), True, False, False, False))
        elif m.group(2) is not None:
            segs.append((m.group(2), False, True, False, False))
        elif m.group(3) is not None:
            segs.append((m.group(3), False, False, True, False))
        elif m.group(4) is not None:
            segs.append((m.group(4), False, False, False, True))
        last = m.end()
    if last < len(txt):
        segs.append((txt[last:], False, False, False, False))
    return segs or [(txt, False, False, False, False)]


def _add_inline_runs(p, txt, size, color_hex, base_bold=False, base_italic=False):
    """Add runs with inline markdown formatting (**bold**, *italic*, ~~strike~~, `code`)."""
    sz_val = str(int(size * 100))
    clr = color_hex.upper()
    for (seg, bold, italic, strike, code) in _parse_inline(txt):
        if not seg:
            continue
        eff_bold = base_bold or bold
        eff_italic = base_italic or italic
        for is_jp, chunk in _split_script(seg):
            r = _el("r")
            rPr_attrib = {
                "lang": "ja-JP" if is_jp else "en-US",
                "sz": sz_val,
                "b": "1" if eff_bold else "0",
                "i": "1" if eff_italic else "0",
                "dirty": "0",
            }
            if strike:
                rPr_attrib["strike"] = "sngStrike"
            rPr = _el("rPr", rPr_attrib)
            fill = _el("solidFill")
            fill.append(_el("srgbClr", {"val": clr}))
            rPr.append(fill)
            typeface = "Meiryo" if is_jp else ("Consolas" if code else "Segoe UI")
            rPr.append(_el("ea" if is_jp else "latin", {"typeface": typeface}))
            t = _el("t")
            t.text = chunk
            r.append(rPr); r.append(t)
            p._p.append(r)


def parse_md(md):
    sections, cur_title, cur_items = [], "", []
    ordered_idx = 0
    for line in md.strip().split("\n"):
        s = line.rstrip()
        stripped = s.strip()
        if not stripped:
            ordered_idx = 0
            continue
        indent = len(s) - len(s.lstrip(" \t"))
        m = re.match(r"^#{1,6}\s+(.*)", stripped)
        if m:
            if cur_title or cur_items:
                sections.append((cur_title, cur_items))
                cur_items = []
            cur_title = m.group(1)
            ordered_idx = 0
        elif re.match(r"^>\s*", stripped):
            cur_items.append(("blockquote", re.sub(r"^>\s*", "", stripped)))
            ordered_idx = 0
        elif re.match(r"^[-*+]\s+", stripped):
            content = re.sub(r"^[-*+]\s+", "", stripped)
            cur_items.append(("bullet2" if indent >= 2 else "bullet", content))
            ordered_idx = 0
        elif re.match(r"^\d+\.\s+", stripped):
            content = re.sub(r"^\d+\.\s+", "", stripped)
            if indent >= 2:
                cur_items.append(("ordered2", content))
            else:
                ordered_idx += 1
                cur_items.append(("ordered", f"{ordered_idx}. {content}"))
        else:
            cur_items.append(("para", stripped))
            ordered_idx = 0
    if cur_title or cur_items:
        sections.append((cur_title, cur_items))
    return sections


def _set_table_cell(cell, cell_text, size, bold, text_color, bg_color, border_color):
    from pptx.oxml import parse_xml as _px
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tf = cell.text_frame
    tf.margin_left = tf.margin_right = Pt(7)
    tf.margin_top  = tf.margin_bottom = Pt(3)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_inline_runs(tf.paragraphs[0], cell_text, size, text_color, base_bold=bold)
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    tcPr.insert(0, _px(
        f'<a:solidFill xmlns:a="{NS_A}"><a:srgbClr val="{bg_color.upper()}"/></a:solidFill>'
    ))
    border_emu = str(int(0.75 * 12700))
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        for old in tcPr.findall(qn(f"a:{edge}")):
            tcPr.remove(old)
        tcPr.append(_px(
            f'<a:{edge} xmlns:a="{NS_A}" w="{border_emu}" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{border_color.upper()}"/></a:solidFill>'
            f'<a:prstDash val="solid"/></a:{edge}>'
        ))


# ── グリッド計算 ────────────────────────────────────────

def compute_cells(sd):
    """grid定義からセル位置を計算（JSエンジンと同一公式）"""
    mainX  = L["mainPadX"]
    mainY  = L["headerH"] + L["mainPadY"]
    mainW  = L["slideW"]  - L["mainPadX"] * 2
    mainH  = L["footerY"] - mainY - L["mainPadY"]
    gap    = L["gridGap"]
    grid   = sd["grid"]
    numRows = len(grid)
    numCols = max(sum(c.get("span", 1) for c in row) for row in grid)
    ratios  = sd.get("rowHeightRatios", [1.0 / numRows] * numRows)
    rowHeights = [mainH * r for r in ratios]

    def rowTop(r):
        return mainY + sum(rowHeights[:r])

    unitW = (mainW - gap * (numCols + 1)) / numCols
    cells = []
    for ri, row in enumerate(grid):
        col = 0
        for cell in row:
            span = cell.get("span", 1)
            cx = mainX + gap + col * (unitW + gap)
            cy = rowTop(ri) + (gap if ri == 0 else gap / 2)
            cw = unitW * span + gap * (span - 1)
            ch = rowHeights[ri] - (gap if ri == 0 else gap / 2) \
                                - (gap if ri == numRows - 1 else gap / 2)
            cells.append({"cell": cell, "x": cx, "y": cy, "w": cw, "h": ch})
            col += span
    return cells


# ── セル描画 ────────────────────────────────────────────

def render_cell(slide, ci):
    cell = ci["cell"]
    cx, cy, cw, ch = ci["x"], ci["y"], ci["w"], ci["h"]
    t    = cell.get("type", "card")
    padX = L["cardPadX"]; padY = L["cardPadY"]

    if t == "section":
        rect(slide, cx, cy, cw, ch, C.get("bgBox", "E8EDF2"), C["border"], 0.75)
        md = cell.get("markdown", "")
        if md:
            sections  = parse_md(md)
            all_items = []
            for title, items in sections:
                if title:
                    all_items.append(("heading", title))
                all_items.extend(items)
            if all_items:
                md_content(slide, all_items,
                           cx + padX, cy + padY, cw - padX * 2, ch - padY * 2,
                           F["bgBody"]["size"], C["text"],
                           head_size=F["bgTitle"]["size"])
            else:
                text(slide, md, cx + padX, cy + padY,
                     cw - padX * 2, ch - padY * 2,
                     F["bgBody"]["size"], False, C["text"], anchor=MSO_ANCHOR.TOP)

    elif t == "card":
        rect(slide, cx, cy, cw, ch, C["surface"], C["border"], 0.75)
        md = cell.get("markdown", "")
        sections = parse_md(md)
        card_title, items = sections[0] if sections else ("", [])
        divY_abs = cy + L["cardDivY"]
        if card_title:
            text(slide, card_title,
                 cx + padX, cy + padY, cw - padX * 2, 0.35,
                 F["cardTitle"]["size"], F["cardTitle"]["bold"], C["text"],
                 anchor=MSO_ANCHOR.TOP)
            hline(slide, cx + padX, divY_abs, cw - padX * 2, C["border"], 1.0)
            if items:
                md_content(slide, items,
                           cx + padX, divY_abs + 0.08, cw - padX * 2,
                           cy + ch - divY_abs - 0.08 - padY,
                           F["cardBody"]["size"], C["text"])
        elif items:
            md_content(slide, items,
                       cx + padX, cy + padY, cw - padX * 2, ch - padY * 2,
                       F["cardBody"]["size"], C["text"])

    elif t == "table":
        head     = cell.get("head", [])
        rows     = cell.get("rows", [])
        has_head = bool(head)
        num_cols = max(len(head) if head else 0,
                       max((len(r) for r in rows), default=0))
        num_rows = len(rows) + (1 if has_head else 0)
        if num_rows > 0 and num_cols > 0:
            tbl_shape = slide.shapes.add_table(
                num_rows, num_cols, Inches(cx), Inches(cy), Inches(cw), Inches(ch))
            tbl = tbl_shape.table
            tbl_pr = tbl._tbl.find(qn("a:tblPr"))
            if tbl_pr is not None:
                for attr in ("firstRow", "bandRow", "firstCol"):
                    tbl_pr.attrib.pop(attr, None)
            for j in range(num_cols):
                tbl.columns[j].width = Inches(cw / num_cols)
            for i in range(num_rows):
                tbl.rows[i].height = Inches(ch / num_rows)
            if has_head:
                for j, ct in enumerate(head[:num_cols]):
                    _set_table_cell(tbl.cell(0, j), ct,
                                    F["tableHead"]["size"], True,
                                    C["tableHeadText"], C["tableHead"], C["tableBorder"])
            row_offset = 1 if has_head else 0
            for i, row in enumerate(rows):
                bg = C["tableRowAlt"] if i % 2 == 1 else C["tableRow"]
                for j, ct in enumerate(row[:num_cols]):
                    if j == 0:
                        _set_table_cell(tbl.cell(i + row_offset, j), ct,
                                        F["tableHead"]["size"], True,
                                        C["tableHeadText"], C["tableHead"], C["tableBorder"])
                    else:
                        _set_table_cell(tbl.cell(i + row_offset, j), ct,
                                        F["tableBody"]["size"], False,
                                        C["tableText"], bg, C["tableBorder"])

    elif t == "conclusion":
        accentW = L.get("conclusionAccentW", 0.07)
        rect(slide, cx, cy, cw, ch, C["conclusionBg"], C["conclusionBorder"], 1.5)
        rect(slide, cx, cy, accentW, ch, C["conclusionBorder"])
        md        = cell.get("markdown", "")
        content_x = cx + accentW + padX
        content_w = cw - accentW - padX * 2
        sections  = parse_md(md) if md else []
        all_items = []
        for title, items in sections:
            if title:
                all_items.append(("heading", title))
            all_items.extend(items)
        if all_items:
            md_content(slide, all_items,
                       content_x, cy + padY, content_w, ch - padY * 2,
                       F["conclBody"]["size"], C["conclusionText"],
                       head_size=F["conclTitle"]["size"])
        elif md:
            text(slide, md, content_x, cy + padY, content_w, ch - padY * 2,
                 F["conclBody"]["size"], False, C["conclusionText"],
                 anchor=MSO_ANCHOR.TOP)

    elif t == "arrow":
        arrowW = cw * 0.55
        arrowH = ch * 0.75
        arrowX = cx + (cw - arrowW) / 2
        arrowY = cy + (ch - arrowH) / 2
        down_triangle_shape(slide, arrowX, arrowY, arrowW, arrowH,
                            C.get("arrow", "6B7A99"))

    elif t == "step_head":
        arrow_shape(slide, cx, cy, cw, ch,
                    C.get("stepFill", "D0D0D0"), C.get("stepBorder", "CCCCCC"))
        label = cell.get("markdown", "")
        if label:
            text(slide, label, cx + 0.15, cy, cw * 0.85, ch,
                 F["stepLabel"]["size"], F["stepLabel"]["bold"],
                 C.get("stepText", "333333"))

    elif t == "image":
        src    = cell.get("src", "")
        label  = cell.get("markdown", "画像 / Image")
        grid_n = int(cell.get("gridN", 3))
        rect(slide, cx, cy, cw, ch, C["surface"], C["border"], 0.75)
        img_path = os.path.join(HERE, src) if src else ""
        if src and os.path.exists(img_path):
            pic = slide.shapes.add_picture(
                img_path, Inches(cx), Inches(cy), Inches(cw), Inches(ch))
            _no_shadow(pic)
        else:
            text(slide, label,
                 cx, cy, cw, ch,
                 F["bodyText"]["size"], False, C["textMuted"],
                 align=PP_ALIGN.CENTER)

    elif t == "plain":
        md        = cell.get("markdown", "")
        sections  = parse_md(md)
        all_items = []
        for title, items in sections:
            if title:
                all_items.append(("heading", title))
            all_items.extend(items)
        hPadX = L["headerPadX"]; hPadY = L["mainPadY"]
        if all_items:
            md_content(slide, all_items,
                       cx + hPadX, cy + hPadY, cw - hPadX * 2, ch - hPadY * 2,
                       F["bodyText"]["size"], C["text"],
                       head_size=F["bodyHead"]["size"])


# ── ヘッダー・フッター・カバー ──────────────────────────

def render_header(slide, sd):
    rect(slide, 0, 0, L["slideW"], L["headerH"], C["headerBg"])
    header = sd.get("header", {})
    if isinstance(header, str):
        header = {"title": header}
    text(slide, header.get("title", ""),
         L["headerPadX"], L["headerPadY"] + 0.013,
         L["slideW"] - L["headerPadX"] * 2, 0.45,
         F["title"]["size"], F["title"]["bold"], C["headerText"])
    text(slide, header.get("message", ""),
         L["headerPadX"], L["headerPadY"] + 0.45 + 0.042,
         L["slideW"] - L["headerPadX"] * 2, 0.3,
         F["message"]["size"], F["message"]["bold"], C["headerText"])


def render_footer(slide, sd):
    footer_cy = L["footerY"] + L["footerH"] / 2
    page_h    = 0.3
    hline(slide, 0, L["footerY"], L["slideW"], C["border"], 0.75)
    box = text(slide, sd.get("page", ""),
               L["footerPadX"], footer_cy - page_h / 2, 2, page_h,
               F["footer"]["size"], False, C["textMuted"])
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    logo_path = os.path.join(HERE, sd.get("logo", "logo.png"))
    if os.path.exists(logo_path):
        pic = slide.shapes.add_picture(
            logo_path,
            Inches(L["slideW"] - L["footerPadX"] - 1.0),
            Inches(footer_cy - L["logoH"] / 2),
            height=Inches(L["logoH"]),
        )
        _no_shadow(pic)


def render_cover(slide, sd):
    bg_path = os.path.join(HERE, sd.get("bg", "background.png"))
    if os.path.exists(bg_path):
        slide.shapes.add_picture(
            bg_path, Inches(0), Inches(0),
            width=Inches(L["slideW"]), height=Inches(L["slideH"])
        )
    cx, cw = 0.5, 5.5
    title_lines = sd.get("title", "").split("\n")
    box = slide.shapes.add_textbox(Inches(cx), Inches(1.36), Inches(cw), Inches(1.35))
    tf  = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _fix_bodyPr(tf)
    for i, line in enumerate(title_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _add_inline_runs(p, line, F["coverTitle"]["size"], C["coverTitle"], base_bold=F["coverTitle"]["bold"])
        _set_line_spacing(p, F["coverTitle"]["size"], mult=1.4)
    _no_shadow(box)
    hline(slide, 0.5, 2.81, 4.0, C["coverDivider"], 0.75)
    for i, key in enumerate(["affiliation", "presenter", "date"]):
        text(slide, sd.get(key, ""), cx, 3.01 + i * 0.35, cw, 0.35,
             F["coverMeta"]["size"], False, C["coverMeta"])


def render_plain2col(slide, sd):
    """plain2col エンジン用（columns[] 構造）"""
    colY   = L["headerH"] + L["mainPadY"]
    colH   = L["footerY"] - colY - L["mainPadY"]
    totalW = L["slideW"] - L["mainPadX"] * 2
    colW   = (totalW - L["gridGap"]) / 2
    col1X  = L["mainPadX"]
    col2X  = L["mainPadX"] + colW + L["gridGap"]
    midX   = col1X + colW + L["gridGap"] / 2
    positions = [(col1X, colY, colW, colH), (col2X, colY, colW, colH)]
    vline(slide, midX, colY, colH, C["border"], 1.0)
    for i, col in enumerate(sd.get("columns", [])):
        if i >= len(positions):
            break
        bx, by, bw, bh = positions[i]
        padX = L["headerPadX"]; padY = L["mainPadY"]
        sections  = parse_md(col.get("markdown", ""))
        all_items = []
        for title, items in sections:
            if title:
                all_items.append(("heading", title))
            all_items.extend(items)
        if all_items:
            md_content(slide, all_items,
                       bx + padX, by + padY, bw - padX * 2, bh - padY * 2,
                       F["bodyText"]["size"], C["text"],
                       head_size=F["bodyHead"]["size"])


# ── メイン ──────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(L["slideW"])
    prs.slide_height = Inches(L["slideH"])

    targets = {k: v for k, v in TEMPLATES.items()
               if _arg is None or k == _arg}
    if not targets:
        print(f"テンプレートが見つかりません: {_arg}"); return

    total = 0
    print(f"[Info] テンプレート数: {len(targets)}")
    print("=" * 60)

    for tmpl_key, tmpl in targets.items():
        slides_data = tmpl.get("SLIDES", [])
        engine      = tmpl.get("engine", "area")
        print(f"[*] {tmpl_key} ({tmpl.get('name', '')})...", end=" ")
        try:
            for sd in slides_data:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                bg = slide.background
                bg.fill.solid()
                bg.fill.fore_color.rgb = rgb(C.get("bg", "FFFFFF"))

                if sd.get("layout") == "cover":
                    render_cover(slide, sd)
                elif engine == "plain2col":
                    render_header(slide, sd)
                    render_plain2col(slide, sd)
                    render_footer(slide, sd)
                else:
                    render_header(slide, sd)
                    if "grid" in sd:
                        for ci in compute_cells(sd):
                            render_cell(slide, ci)
                    render_footer(slide, sd)

                # ノート（発表者原稿）の書き込み
                note_text = sd.get("note", "")
                if note_text:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = note_text

            total += len(slides_data)
            print(f"OK ({len(slides_data)} slides)")
        except Exception as e:
            import traceback
            print(f"ERROR: {e}")
            traceback.print_exc()

    print("=" * 60)
    print(f"[Done] 合計スライド数: {total}")
    out_name = f"{_arg}.pptx" if _arg else "all_templates.pptx"
    out_path = os.path.join(HERE, out_name)
    prs.save(out_path)
    print(f"[Save] {out_path}")
    print(f"[Size] {os.path.getsize(out_path) / 1024:.1f} KB")


if __name__ == "__main__":
    main()
